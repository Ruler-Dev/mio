"""Document-generation skills using the same libraries as Anthropic's Claude
Skills (https://github.com/anthropics/skills):

- generate_docx   → python-docx
- generate_xlsx   → openpyxl
- generate_pptx   → python-pptx
- generate_pdf_report → reportlab (with optional embedded matplotlib charts)

All outputs land in ~/Downloads.
"""

from __future__ import annotations

import io
import time
from pathlib import Path

from mio.webui.safe_files import (
    UnsafePathError,
    downloads_input_path,
    downloads_output_path,
    open_binary_no_follow,
)


def _output_path(filename: str | None, ext: str) -> Path:
    fn = filename or f"mio-{int(time.time())}{ext}"
    return downloads_output_path(fn, ext)


# ---------------------------------------------------------------------------
# Office (docx/xlsx/pptx) preset catalog. Structurally simpler than the PDF
# catalog — each preset is a color+font bundle; Office layouts are driven
# by the library defaults so we can't hand-compose decoration the way we
# can on a canvas. Still gives the user ~30 distinct vibes.
# ---------------------------------------------------------------------------

_OFFICE_PRESETS: dict = {
    # Blues
    "azure": dict(accent="3B82F6", header="1F4E78", row_alt="EEF5FC", text="111111", family="sans"),
    "cobalt": dict(accent="1D4ED8", header="172554", row_alt="EEF2FF", text="111111", family="sans"),
    "sky": dict(accent="0284C7", header="0C4A6E", row_alt="F0F9FF", text="111111", family="sans"),
    "ocean": dict(accent="0EA5E9", header="075985", row_alt="E0F2FE", text="111111", family="sans"),
    "navy": dict(accent="0F2A5A", header="0F2A5A", row_alt="EEF1F8", text="111111", family="serif"),
    "oxford": dict(accent="1E3A8A", header="172554", row_alt="EFF6FF", text="111111", family="serif"),
    # Slate / mono
    "slate": dict(accent="475569", header="1E293B", row_alt="F1F5F9", text="111111", family="sans"),
    "graphite": dict(accent="1F2937", header="111827", row_alt="F3F4F6", text="111111", family="sans"),
    "minimal": dict(accent="111111", header="111111", row_alt="F5F5F5", text="111111", family="sans"),
    "bone": dict(accent="1A1A1A", header="262626", row_alt="F5F5F4", text="111111", family="serif"),
    # Greens
    "emerald": dict(accent="059669", header="064E3B", row_alt="ECFDF5", text="111111", family="sans"),
    "forest": dict(accent="15803D", header="14532D", row_alt="F0FDF4", text="111111", family="sans"),
    "sage": dict(accent="4D7C5B", header="365240", row_alt="F3F6F3", text="111111", family="sans"),
    "teal": dict(accent="0D9488", header="134E4A", row_alt="F0FDFA", text="111111", family="sans"),
    # Reds / warms
    "crimson": dict(accent="B91C1C", header="7F1D1D", row_alt="FEF2F2", text="111111", family="serif"),
    "ruby": dict(accent="DC2626", header="991B1B", row_alt="FEF2F2", text="111111", family="sans"),
    "rose": dict(accent="E11D48", header="881337", row_alt="FFF1F2", text="111111", family="sans"),
    "amber": dict(accent="D97706", header="78350F", row_alt="FFFBEB", text="111111", family="sans"),
    "sunset": dict(accent="EA580C", header="9A3412", row_alt="FEF3E8", text="111111", family="sans"),
    "terracotta": dict(accent="C2410C", header="7C2D12", row_alt="FFEDD5", text="111111", family="serif"),
    # Purples
    "violet": dict(accent="7C3AED", header="4C1D95", row_alt="F5F3FF", text="111111", family="sans"),
    "indigo": dict(accent="4F46E5", header="312E81", row_alt="EEF2FF", text="111111", family="sans"),
    "plum": dict(accent="7E22CE", header="581C87", row_alt="FAF5FF", text="111111", family="serif"),
    # Corporate / heritage
    "corporate": dict(accent="0F2A5A", header="0F2A5A", row_alt="EEF1F8", text="111111", family="serif"),
    "embassy": dict(accent="7F1D1D", header="450A0A", row_alt="FEF2F2", text="111111", family="serif"),
    "walnut": dict(accent="78350F", header="422006", row_alt="FEF3E8", text="111111", family="serif"),
    "brass": dict(accent="A16207", header="713F12", row_alt="FEFCE8", text="111111", family="serif"),
    "cambridge": dict(accent="166534", header="14532D", row_alt="F0FDF4", text="111111", family="serif"),
    "academic": dict(accent="374151", header="374151", row_alt="F3F4F6", text="111111", family="serif"),
    # Playful
    "neon": dict(accent="EC4899", header="831843", row_alt="FDF2F8", text="111111", family="sans"),
    "tropical": dict(accent="06B6D4", header="155E75", row_alt="ECFEFF", text="111111", family="sans"),
    "bubblegum": dict(accent="EC4899", header="BE185D", row_alt="FDF2F8", text="111111", family="sans"),
    "citrus": dict(accent="84CC16", header="365314", row_alt="F7FEE7", text="111111", family="sans"),
    "peach": dict(accent="FB923C", header="9A3412", row_alt="FFEDD5", text="111111", family="sans"),
    # Technical / dark-accent on white (no true dark background because
    # Office apps default to white pages; these just use darker accents)
    "carbon": dict(accent="27272A", header="18181B", row_alt="F4F4F5", text="111111", family="mono"),
    "midnight": dict(accent="1E293B", header="0F172A", row_alt="F1F5F9", text="111111", family="sans"),
    "obsidian": dict(accent="111827", header="030712", row_alt="F3F4F6", text="111111", family="sans"),
    "blueprint": dict(accent="0369A1", header="0C4A6E", row_alt="E0F2FE", text="111111", family="mono"),
    "terminal": dict(accent="166534", header="14532D", row_alt="F0FDF4", text="111111", family="mono"),
    # Editorial
    "editorial": dict(accent="B91C1C", header="1A1A1A", row_alt="FBF3EE", text="111111", family="serif"),
    "vogue": dict(accent="000000", header="262626", row_alt="FAFAFA", text="111111", family="serif"),
    "journal": dict(accent="374151", header="1F2937", row_alt="F9FAFB", text="111111", family="serif"),
    "parchment": dict(accent="78350F", header="451A03", row_alt="FEF3E8", text="111111", family="serif"),
}

_OFFICE_FONTS = {
    "sans": dict(heading="Calibri", body="Calibri"),
    "serif": dict(heading="Cambria", body="Cambria"),
    "mono": dict(heading="Consolas", body="Consolas"),
}

# Per-format pools so decks don't pick serif-heavy presets by default and
# spreadsheets lean on the clearer sans palettes.
_OFFICE_POOLS: dict = {
    "docx": list(_OFFICE_PRESETS.keys()),
    "xlsx": [
        "azure",
        "cobalt",
        "sky",
        "ocean",
        "navy",
        "oxford",
        "slate",
        "graphite",
        "minimal",
        "emerald",
        "forest",
        "teal",
        "indigo",
        "violet",
        "corporate",
        "academic",
        "tropical",
        "crimson",
        "amber",
        "sunset",
        "rose",
        "carbon",
        "midnight",
        "blueprint",
        "terminal",
    ],
    "pptx": [
        "azure",
        "cobalt",
        "ocean",
        "navy",
        "slate",
        "graphite",
        "minimal",
        "emerald",
        "teal",
        "indigo",
        "violet",
        "crimson",
        "ruby",
        "rose",
        "sunset",
        "amber",
        "neon",
        "tropical",
        "bubblegum",
        "citrus",
        "peach",
        "plum",
        "corporate",
        "embassy",
        "walnut",
        "carbon",
        "midnight",
        "obsidian",
        "editorial",
        "vogue",
        "journal",
        "blueprint",
    ],
}


def _pick_office_preset(fmt: str, *text_hints: str) -> str:
    """Auto-pick an Office preset. Reuses the PDF keyword table but with
    fallbacks appropriate for Office (no decoration variants)."""
    import re as _re

    haystack = " ".join(t for t in text_hints if t).lower()
    tokens = set(_re.findall(r"[a-z0-9]+", haystack)) if haystack else set()
    pool = _OFFICE_POOLS.get(fmt, _OFFICE_POOLS["docx"])
    if tokens:
        for keywords, preset in _PDF_KEYWORD_HINTS:

            def matches(k: str) -> bool:
                return (k in haystack) if " " in k else (k in tokens)

            if any(matches(k) for k in keywords):
                if preset in pool:
                    return preset
                if preset in _OFFICE_PRESETS:
                    return preset
    import hashlib as _hl

    h = int(_hl.sha1(haystack.encode("utf-8")).hexdigest()[:8], 16) if haystack else 0
    return pool[h % len(pool)]


def _office_theme(
    name: str,
    fmt: str,
    *hints: str,
    color: str | None = None,
    accent_color: str | None = None,
    text_color: str | None = None,
    header_color: str | None = None,
    background_color: str | None = None,  # reserved for future DOCX/PPTX fill support
) -> dict:
    """Resolve an Office preset. `color="emerald"` overrides only the color
    triple (accent / header / row_alt). Surgical overrides (`accent_color`,
    `text_color`, `header_color`) win over both preset and palette; they
    accept natural names or hex.
    """
    key = (name or "auto").strip().lower()
    if key in ("auto", "", "default"):
        key = _pick_office_preset(fmt, *hints)
    p = _OFFICE_PRESETS.get(key) or _OFFICE_PRESETS["azure"]
    f = _OFFICE_FONTS.get(p.get("family", "sans"))
    accent = p["accent"]
    header = p["header"]
    row_alt = p["row_alt"]
    text = p.get("text", "111111")
    override = _resolve_color_palette(color)
    if override:
        # Office palette is stored without "#" prefix — strip it
        accent = override["accent"].lstrip("#").upper()
        header = override["head_bg"].lstrip("#").upper()
        row_alt = override["row_alt"].lstrip("#").upper()

    def _office_hex(v: str | None) -> str | None:
        h = _coerce_color(v)
        return h.lstrip("#").upper() if h else None

    if h := _office_hex(accent_color):
        accent = h
    if h := _office_hex(header_color):
        header = h
    if h := _office_hex(text_color):
        text = h

    return {
        "name": key,
        "color_override": (color or "").lower() if override else None,
        "accent": accent,
        "header": header,
        "row_alt": row_alt,
        "text": text,
        "heading_font": f["heading"],
        "body_font": f["body"],
    }


OFFICE_PRESETS = tuple(sorted(_OFFICE_PRESETS.keys()))


# ============================================================
# DOCX — python-docx
# ============================================================
def _docx_add_inline(p, text: str) -> None:
    """Append text to a python-docx paragraph, honoring **bold** and `code`
    inline runs. Keeps rendering simple but avoids raw markdown showing up."""
    import re as _re

    tokens = _re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text)
    for tok in tokens:
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**"):
            run = p.add_run(tok[2:-2])
            run.bold = True
        elif tok.startswith("`") and tok.endswith("`"):
            run = p.add_run(tok[1:-1])
            run.font.name = "Menlo"
        else:
            p.add_run(tok)


def generate_docx(
    title: str,
    content: str,
    filename: str | None = None,
    author: str = "",
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Build a Word document from markdown content.

    Uses the same markdown parser as generate_pdf_report — supports
    ATX headings (# / ## / ###), paragraphs, bulleted lists, and GFM pipe
    tables. Inline **bold** and `code` are rendered.
    """
    try:
        from docx import Document
        from docx.shared import RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
    except ImportError:
        return {"skill": "generate_docx", "error": "python-docx not installed."}

    T = _office_theme(
        preset,
        "docx",
        title,
        (content or "")[:400],
        color=color,
        accent_color=accent_color,
        text_color=text_color,
    )

    def _hex_to_rgb(h):
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    ACCENT = _hex_to_rgb(T["accent"])
    HEADER = _hex_to_rgb(T["header"])

    doc = Document()
    if author:
        doc.core_properties.author = author
    doc.core_properties.title = title

    # Preset font on the default style so headings + body pick up the vibe
    try:
        style = doc.styles["Normal"]
        style.font.name = T["body_font"]
    except KeyError:
        pass

    # Title — colored to match the preset accent
    t = doc.add_heading(title, level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for r in t.runs:
        r.font.color.rgb = HEADER
        r.font.name = T["heading_font"]

    blocks = _markdown_to_sections(content or "")
    for b in blocks:
        kind = b.get("kind")
        if kind == "heading":
            lvl = int(b.get("level", 1))
            h = doc.add_heading(b.get("text", ""), level=min(max(lvl, 1), 3))
            for r in h.runs:
                r.font.color.rgb = ACCENT if lvl >= 2 else HEADER
                r.font.name = T["heading_font"]
        elif kind == "paragraph":
            p = doc.add_paragraph()
            _docx_add_inline(p, b.get("text", ""))
        elif kind == "bullets":
            for it in b.get("items", []):
                p = doc.add_paragraph(style="List Bullet")
                _docx_add_inline(p, str(it))
        elif kind == "table":
            headers = b.get("headers") or []
            rows = b.get("rows") or []
            if not headers and not rows:
                continue
            ncols = len(headers) if headers else max((len(r) for r in rows), default=1)
            tbl = doc.add_table(rows=1 + len(rows), cols=ncols)
            tbl.style = "Light Grid Accent 1"
            tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
            # Header row
            for j, h in enumerate(headers[:ncols]):
                cell = tbl.rows[0].cells[j]
                cell.text = str(h)
                for r in cell.paragraphs[0].runs:
                    r.bold = True
            # Data rows
            for i, row in enumerate(rows, start=1):
                row = list(row) if isinstance(row, (list, tuple)) else [row]
                for j in range(ncols):
                    tbl.rows[i].cells[j].text = str(row[j]) if j < len(row) else ""
            doc.add_paragraph("")  # breathing room after table

    out = _output_path(filename, ".docx")
    doc.save(str(out))
    return {"skill": "generate_docx", "path": str(out), "filename": out.name}


# ============================================================
# XLSX — openpyxl
# ============================================================
def generate_xlsx(
    title: str,
    headers: list[str],
    rows: list[list],
    filename: str | None = None,
    sheet_name: str = "Sheet1",
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Build an Excel workbook with one sheet of tabular data."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        return {"skill": "generate_xlsx", "error": "openpyxl not installed."}

    if not headers:
        return {"skill": "generate_xlsx", "error": "headers required"}

    T = _office_theme(
        preset,
        "xlsx",
        title,
        sheet_name,
        " ".join(str(h) for h in headers[:8]),
        color=color,
        accent_color=accent_color,
        text_color=text_color,
    )

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name[:31] or "Sheet1"

    # Optional title row
    if title:
        tcell = ws.cell(row=1, column=1, value=title)
        tcell.font = Font(bold=True, size=14, color=T["header"], name=T["heading_font"])
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
        header_row = 3
    else:
        header_row = 1

    # Header styling driven by preset
    header_fill = PatternFill("solid", fgColor=T["header"])
    alt_fill = PatternFill("solid", fgColor=T["row_alt"])
    thin = Side(border_style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for col_idx, h in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=col_idx, value=str(h))
        cell.font = Font(bold=True, color="FFFFFF", name=T["heading_font"])
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = border

    # Data rows — zebra-stripe with the preset's row_alt color
    for r_idx, row in enumerate(rows, start=header_row + 1):
        is_alt = (r_idx - header_row) % 2 == 0
        for c_idx, val in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx, value=val)
            cell.border = border
            cell.font = Font(name=T["body_font"])
            if is_alt:
                cell.fill = alt_fill

    # Auto-size columns (best effort)
    for col_idx, h in enumerate(headers, start=1):
        max_len = len(str(h))
        for r in rows:
            if col_idx - 1 < len(r):
                max_len = max(max_len, len(str(r[col_idx - 1])))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 2, 40)

    ws.freeze_panes = ws.cell(row=header_row + 1, column=1)

    out = _output_path(filename, ".xlsx")
    wb.save(str(out))
    return {
        "skill": "generate_xlsx",
        "path": str(out),
        "filename": out.name,
        "rows": len(rows),
        "cols": len(headers),
    }


# ============================================================
# PPTX — python-pptx
# ============================================================
def generate_pptx(
    title: str,
    slides: list[dict],
    filename: str | None = None,
    subtitle: str = "",
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Build a PowerPoint deck.

    `slides` is a list of dicts; each dict may contain:
      - 'title': slide title
      - 'content': markdown-ish bullets (one per line starting with '- ')
      - 'notes':   speaker notes
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return {"skill": "generate_pptx", "error": "python-pptx not installed."}

    hints = " ".join([title, subtitle or ""] + [s.get("title", "") for s in (slides or [])[:6]])
    T = _office_theme(
        preset,
        "pptx",
        hints,
        color=color,
        accent_color=accent_color,
        text_color=text_color,
    )

    def _hex_to_rgb(h):
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    ACCENT = _hex_to_rgb(T["accent"])
    HEADER = _hex_to_rgb(T["header"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _decorate(slide, with_title_color=True):
        """Add a skinny accent bar along the left edge of the slide."""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        # Optional footer dot-bar
        dot = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.18), prs.slide_width, Inches(0.18)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = HEADER
        dot.line.fill.background()
        if with_title_color and slide.shapes.title:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = HEADER
                    run.font.name = T["heading_font"]
                    run.font.bold = True

    # Title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        try:
            title_slide.placeholders[1].text = subtitle or ""
        except Exception:
            pass
    _decorate(title_slide)

    # Content slides
    content_layout = prs.slide_layouts[1]
    for s in slides:
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("title", "")
        _decorate(slide)
        # body placeholder
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is not None:
            tf = body.text_frame
            tf.clear()
            lines = (s.get("content", "") or "").split("\n")
            first = True
            for line in lines:
                stripped = line.strip()
                if not stripped:
                    continue
                if stripped.startswith("- ") or stripped.startswith("* "):
                    stripped = stripped[2:]
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = stripped
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(20)
                first = False
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    out = _output_path(filename, ".pptx")
    prs.save(str(out))
    return {
        "skill": "generate_pptx",
        "path": str(out),
        "filename": out.name,
        "slide_count": len(prs.slides),
    }


# ============================================================
# PDF REPORT — reportlab (with optional matplotlib charts)
# ============================================================
def _markdown_to_sections(md: str) -> list[dict]:
    """Very lightweight markdown → section-blocks converter.

    Supports: ATX headings (# / ## / ###), paragraphs (blank-line separated),
    bulleted lists (-, *, +), and GitHub-flavoured pipe tables.
    Emphasis markers (**, *, `) are left as plain text (reportlab renders
    limited inline markup; we keep it simple and readable).
    """
    import re as _re

    lines = [ln.rstrip() for ln in (md or "").splitlines()]
    out: list[dict] = []
    i = 0
    para: list[str] = []

    def flush_para():
        if para:
            txt = " ".join(s.strip() for s in para if s.strip())
            if txt:
                out.append({"kind": "paragraph", "text": txt})
            para.clear()

    while i < len(lines):
        ln = lines[i]
        if not ln.strip():
            flush_para()
            i += 1
            continue
        m = _re.match(r"^(#{1,3})\s+(.*)$", ln)
        if m:
            flush_para()
            out.append({"kind": "heading", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1
            continue
        if _re.match(r"^\s*[-*+]\s+", ln):
            flush_para()
            items = []
            while i < len(lines) and _re.match(r"^\s*[-*+]\s+", lines[i]):
                items.append(_re.sub(r"^\s*[-*+]\s+", "", lines[i]).strip())
                i += 1
            if items:
                out.append({"kind": "bullets", "items": items})
            continue
        # Pipe table: header | --- | rows
        if "|" in ln and i + 1 < len(lines) and _re.match(r"^[\s|:-]+$", lines[i + 1]):
            flush_para()

            def _split(r):
                cells = [c.strip() for c in r.strip().strip("|").split("|")]
                return cells

            headers = _split(ln)
            i += 2
            rows = []
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                rows.append(_split(lines[i]))
                i += 1
            out.append({"kind": "table", "headers": headers, "rows": rows})
            continue
        para.append(ln)
        i += 1
    flush_para()
    return out


# ---------------------------------------------------------------------------
# PDF preset catalog. Each preset bundles colors, typography, and a page-
# decoration recipe. `_pdf_theme("auto", skill=..., content=...)` picks a
# preset automatically based on skill family + content keywords, so the
# model does NOT have to name one — it just calls the skill and gets a
# sensible look; the user can ask for a different vibe later and we'll
# steer toward a different preset.
# ---------------------------------------------------------------------------

_PDF_PRESETS: dict = {
    # Core sans-serif + blue family ——————————————————————————————————
    "azure": dict(accent="#3b82f6", head_bg="#1F4E78", soft="#f4f7fb", decoration="bar", family="modern_sans"),
    "azure-rule": dict(accent="#3b82f6", head_bg="#1e3a8a", soft="#eff6ff", decoration="rule", family="modern_sans"),
    "cobalt": dict(accent="#1d4ed8", head_bg="#172554", soft="#eef2ff", decoration="bar", family="modern_sans"),
    "sky": dict(accent="#0284c7", head_bg="#0c4a6e", soft="#f0f9ff", decoration="bar", family="modern_sans"),
    "ocean": dict(accent="#0ea5e9", head_bg="#075985", soft="#e0f2fe", decoration="block", family="modern_sans"),
    "slate": dict(accent="#475569", head_bg="#1e293b", soft="#f1f5f9", decoration="rule", family="modern_sans"),
    "steel": dict(accent="#334155", head_bg="#0f172a", soft="#f8fafc", decoration="bar", family="modern_sans"),
    "graphite": dict(accent="#1f2937", head_bg="#111827", soft="#f3f4f6", decoration="rule", family="modern_sans"),
    # Green / teal / emerald ——————————————————————————————————————————
    "emerald": dict(accent="#059669", head_bg="#064e3b", soft="#ecfdf5", decoration="bar", family="modern_sans"),
    "forest": dict(accent="#15803d", head_bg="#14532d", soft="#f0fdf4", decoration="rule", family="modern_sans"),
    "sage": dict(accent="#4d7c5b", head_bg="#365240", soft="#f3f6f3", decoration="rule", family="modern_sans"),
    "teal": dict(accent="#0d9488", head_bg="#134e4a", soft="#f0fdfa", decoration="bar", family="modern_sans"),
    "mint": dict(accent="#10b981", head_bg="#065f46", soft="#ecfdf5", decoration="block", family="modern_sans"),
    # Red / orange / warm ————————————————————————————————————————————
    "crimson": dict(accent="#b91c1c", head_bg="#7f1d1d", soft="#fef2f2", decoration="bar", family="classic_serif"),
    "ruby": dict(accent="#dc2626", head_bg="#991b1b", soft="#fef2f2", decoration="rule", family="modern_sans"),
    "rose": dict(accent="#e11d48", head_bg="#881337", soft="#fff1f2", decoration="bar", family="modern_sans"),
    "magenta": dict(accent="#c026d3", head_bg="#701a75", soft="#fdf4ff", decoration="bar", family="modern_sans"),
    "amber": dict(accent="#d97706", head_bg="#78350f", soft="#fffbeb", decoration="bar", family="modern_sans"),
    "sunset": dict(accent="#ea580c", head_bg="#9a3412", soft="#fef3e8", decoration="bar", family="modern_sans"),
    "ember": dict(accent="#dc2626", head_bg="#7c2d12", soft="#fef3e8", decoration="block", family="modern_sans"),
    "terracotta": dict(accent="#c2410c", head_bg="#7c2d12", soft="#ffedd5", decoration="rule", family="classic_serif"),
    "marigold": dict(accent="#f59e0b", head_bg="#78350f", soft="#fffbeb", decoration="bar", family="modern_sans"),
    # Purple / indigo ————————————————————————————————————————————————
    "violet": dict(accent="#7c3aed", head_bg="#4c1d95", soft="#f5f3ff", decoration="bar", family="modern_sans"),
    "indigo": dict(accent="#4f46e5", head_bg="#312e81", soft="#eef2ff", decoration="bar", family="modern_sans"),
    "lavender": dict(accent="#8b5cf6", head_bg="#5b21b6", soft="#f5f3ff", decoration="rule", family="modern_sans"),
    "plum": dict(accent="#7e22ce", head_bg="#581c87", soft="#faf5ff", decoration="block", family="classic_serif"),
    "aubergine": dict(accent="#6b21a8", head_bg="#3b0764", soft="#f5f3ff", decoration="bar", family="classic_serif"),
    # Dark / night ———————————————————————————————————————————————————
    "midnight": dict(
        accent="#60a5fa",
        head_bg="#111827",
        soft="#1f2233",
        decoration="block",
        family="modern_sans",
        page_bg="#0b0f1a",
        text="#f5f5f5",
        muted="#9ca3af",
    ),
    "obsidian": dict(
        accent="#f59e0b",
        head_bg="#0f172a",
        soft="#1e293b",
        decoration="block",
        family="modern_sans",
        page_bg="#020617",
        text="#f1f5f9",
        muted="#94a3b8",
    ),
    "carbon": dict(
        accent="#a78bfa",
        head_bg="#18181b",
        soft="#27272a",
        decoration="rule",
        family="modern_sans",
        page_bg="#09090b",
        text="#fafafa",
        muted="#a1a1aa",
    ),
    "nocturne": dict(
        accent="#22d3ee",
        head_bg="#082f49",
        soft="#0c4a6e",
        decoration="block",
        family="modern_sans",
        page_bg="#082f49",
        text="#e0f2fe",
        muted="#7dd3fc",
    ),
    "eclipse": dict(
        accent="#fb7185",
        head_bg="#1c1917",
        soft="#292524",
        decoration="block",
        family="classic_serif",
        page_bg="#0c0a09",
        text="#fafaf9",
        muted="#a8a29e",
    ),
    # Minimal / mono ——————————————————————————————————————————————————
    "minimal": dict(accent="#111111", head_bg="#111111", soft="#fafafa", decoration="rule", family="modern_sans"),
    "bone": dict(accent="#1a1a1a", head_bg="#262626", soft="#f5f5f4", decoration="none", family="classic_serif"),
    "duotone": dict(accent="#525252", head_bg="#262626", soft="#f5f5f5", decoration="rule", family="modern_sans"),
    "whisper": dict(accent="#737373", head_bg="#404040", soft="#fafafa", decoration="none", family="classic_serif"),
    # Serif / classic / corporate ————————————————————————————————————
    "navy": dict(accent="#0f2a5a", head_bg="#0f2a5a", soft="#eef1f8", decoration="bar", family="classic_serif"),
    "claret": dict(accent="#881337", head_bg="#4c0519", soft="#fff1f2", decoration="bar", family="classic_serif"),
    "oxford": dict(accent="#1e3a8a", head_bg="#172554", soft="#eff6ff", decoration="rule", family="classic_serif"),
    "walnut": dict(accent="#78350f", head_bg="#422006", soft="#fef3e8", decoration="rule", family="classic_serif"),
    "cambridge": dict(accent="#166534", head_bg="#14532d", soft="#f0fdf4", decoration="rule", family="classic_serif"),
    "brass": dict(accent="#a16207", head_bg="#713f12", soft="#fefce8", decoration="bar", family="classic_serif"),
    "parchment": dict(
        accent="#78350f",
        head_bg="#451a03",
        soft="#fef3e8",
        decoration="none",
        family="classic_serif",
        page_bg="#fdf6e3",
    ),
    # Editorial / magazine ————————————————————————————————————————————
    "editorial": dict(accent="#b91c1c", head_bg="#1a1a1a", soft="#fbf3ee", decoration="rule", family="classic_serif"),
    "vogue": dict(accent="#000000", head_bg="#262626", soft="#fafafa", decoration="rule", family="display_serif"),
    "broadsheet": dict(accent="#292524", head_bg="#1c1917", soft="#fafaf9", decoration="rule", family="display_serif"),
    "journal": dict(accent="#374151", head_bg="#1f2937", soft="#f9fafb", decoration="none", family="classic_serif"),
    "manuscript": dict(accent="#525252", head_bg="#262626", soft="#fafaf9", decoration="rule", family="classic_serif"),
    # Monospaced / technical —————————————————————————————————————————
    "terminal": dict(accent="#22c55e", head_bg="#14532d", soft="#f0fdf4", decoration="rule", family="mono"),
    "blueprint": dict(accent="#0284c7", head_bg="#0c4a6e", soft="#e0f2fe", decoration="rule", family="mono"),
    "lab": dict(accent="#4338ca", head_bg="#312e81", soft="#eef2ff", decoration="bar", family="mono"),
    "console": dict(accent="#d97706", head_bg="#78350f", soft="#fffbeb", decoration="rule", family="mono"),
    # Playful / events ————————————————————————————————————————————————
    "neon": dict(accent="#ec4899", head_bg="#831843", soft="#fdf2f8", decoration="bar", family="modern_sans"),
    "tropical": dict(accent="#06b6d4", head_bg="#155e75", soft="#ecfeff", decoration="bar", family="modern_sans"),
    "sunrise": dict(accent="#f97316", head_bg="#9a3412", soft="#ffedd5", decoration="block", family="modern_sans"),
    "bubblegum": dict(accent="#ec4899", head_bg="#be185d", soft="#fdf2f8", decoration="bar", family="modern_sans"),
    "pastel": dict(accent="#a78bfa", head_bg="#6d28d9", soft="#f5f3ff", decoration="rule", family="modern_sans"),
    "citrus": dict(accent="#84cc16", head_bg="#365314", soft="#f7fee7", decoration="bar", family="modern_sans"),
    "peach": dict(accent="#fb923c", head_bg="#9a3412", soft="#ffedd5", decoration="rule", family="modern_sans"),
    # Formal / heritage ——————————————————————————————————————————————
    "embassy": dict(accent="#7f1d1d", head_bg="#450a0a", soft="#fef2f2", decoration="bar", family="classic_serif"),
    "regal": dict(accent="#713f12", head_bg="#422006", soft="#fef3e8", decoration="bar", family="display_serif"),
    "heritage": dict(accent="#064e3b", head_bg="#022c22", soft="#ecfdf5", decoration="bar", family="classic_serif"),
    "academic": dict(accent="#374151", head_bg="#374151", soft="#f3f4f6", decoration="none", family="classic_serif"),
    "corporate": dict(accent="#0f2a5a", head_bg="#0f2a5a", soft="#eef1f8", decoration="bar", family="classic_serif"),
}

_PDF_FONT_STACKS = {
    "modern_sans": dict(heading="Helvetica-Bold", body="Helvetica"),
    "classic_serif": dict(heading="Times-Bold", body="Times-Roman"),
    "display_serif": dict(heading="Times-Bold", body="Times-Roman"),
    "mono": dict(heading="Courier-Bold", body="Courier"),
}


# ---------------------------------------------------------------------------
# Color palettes — 36 named colors. Any palette can be combined with any
# preset layout (decoration + fonts). Lets the user say "same as before
# but in emerald" and only the color triple (accent/header/soft/row_alt)
# is swapped. This is what makes the system feel like Claude: one layout
# × many hues × many decorations = hundreds of real variations.
# ---------------------------------------------------------------------------
_COLOR_PALETTES: dict = {
    # Blues
    "azure": dict(accent="#3b82f6", head_bg="#1F4E78", soft="#f4f7fb", row_alt="EEF5FC"),
    "cobalt": dict(accent="#1d4ed8", head_bg="#172554", soft="#eef2ff", row_alt="EEF2FF"),
    "sky": dict(accent="#0284c7", head_bg="#0c4a6e", soft="#f0f9ff", row_alt="F0F9FF"),
    "ocean": dict(accent="#0ea5e9", head_bg="#075985", soft="#e0f2fe", row_alt="E0F2FE"),
    "navy": dict(accent="#0f2a5a", head_bg="#0f2a5a", soft="#eef1f8", row_alt="EEF1F8"),
    "oxford": dict(accent="#1e3a8a", head_bg="#172554", soft="#eff6ff", row_alt="EFF6FF"),
    "indigo": dict(accent="#4f46e5", head_bg="#312e81", soft="#eef2ff", row_alt="EEF2FF"),
    # Greens / teals
    "emerald": dict(accent="#059669", head_bg="#064e3b", soft="#ecfdf5", row_alt="ECFDF5"),
    "forest": dict(accent="#15803d", head_bg="#14532d", soft="#f0fdf4", row_alt="F0FDF4"),
    "sage": dict(accent="#4d7c5b", head_bg="#365240", soft="#f3f6f3", row_alt="F3F6F3"),
    "teal": dict(accent="#0d9488", head_bg="#134e4a", soft="#f0fdfa", row_alt="F0FDFA"),
    "mint": dict(accent="#10b981", head_bg="#065f46", soft="#ecfdf5", row_alt="ECFDF5"),
    "olive": dict(accent="#65a30d", head_bg="#365314", soft="#f7fee7", row_alt="F7FEE7"),
    # Reds / roses
    "crimson": dict(accent="#b91c1c", head_bg="#7f1d1d", soft="#fef2f2", row_alt="FEF2F2"),
    "ruby": dict(accent="#dc2626", head_bg="#991b1b", soft="#fef2f2", row_alt="FEF2F2"),
    "rose": dict(accent="#e11d48", head_bg="#881337", soft="#fff1f2", row_alt="FFF1F2"),
    "coral": dict(accent="#f43f5e", head_bg="#be123c", soft="#fff1f2", row_alt="FFF1F2"),
    "burgundy": dict(accent="#881337", head_bg="#4c0519", soft="#fff1f2", row_alt="FFF1F2"),
    # Oranges / warms
    "amber": dict(accent="#d97706", head_bg="#78350f", soft="#fffbeb", row_alt="FFFBEB"),
    "sunset": dict(accent="#ea580c", head_bg="#9a3412", soft="#fef3e8", row_alt="FEF3E8"),
    "ember": dict(accent="#dc2626", head_bg="#7c2d12", soft="#fef3e8", row_alt="FEF3E8"),
    "terracotta": dict(accent="#c2410c", head_bg="#7c2d12", soft="#ffedd5", row_alt="FFEDD5"),
    "marigold": dict(accent="#f59e0b", head_bg="#78350f", soft="#fffbeb", row_alt="FFFBEB"),
    "peach": dict(accent="#fb923c", head_bg="#9a3412", soft="#ffedd5", row_alt="FFEDD5"),
    # Purples / magentas
    "violet": dict(accent="#7c3aed", head_bg="#4c1d95", soft="#f5f3ff", row_alt="F5F3FF"),
    "lavender": dict(accent="#8b5cf6", head_bg="#5b21b6", soft="#f5f3ff", row_alt="F5F3FF"),
    "plum": dict(accent="#7e22ce", head_bg="#581c87", soft="#faf5ff", row_alt="FAF5FF"),
    "magenta": dict(accent="#c026d3", head_bg="#701a75", soft="#fdf4ff", row_alt="FDF4FF"),
    "fuchsia": dict(accent="#d946ef", head_bg="#86198f", soft="#fdf4ff", row_alt="FDF4FF"),
    # Neutrals
    "slate": dict(accent="#475569", head_bg="#1e293b", soft="#f1f5f9", row_alt="F1F5F9"),
    "graphite": dict(accent="#1f2937", head_bg="#111827", soft="#f3f4f6", row_alt="F3F4F6"),
    "carbon": dict(accent="#27272a", head_bg="#18181b", soft="#f4f4f5", row_alt="F4F4F5"),
    "midnight": dict(accent="#1e293b", head_bg="#0f172a", soft="#f1f5f9", row_alt="F1F5F9"),
    "mono": dict(accent="#111111", head_bg="#111111", soft="#fafafa", row_alt="F5F5F5"),
    "bone": dict(accent="#1a1a1a", head_bg="#262626", soft="#f5f5f4", row_alt="F5F5F4"),
    # Specials
    "brass": dict(accent="#a16207", head_bg="#713f12", soft="#fefce8", row_alt="FEFCE8"),
    "walnut": dict(accent="#78350f", head_bg="#422006", soft="#fef3e8", row_alt="FEF3E8"),
    "ice": dict(accent="#67e8f9", head_bg="#0e7490", soft="#ecfeff", row_alt="ECFEFF"),
    "neon": dict(accent="#ec4899", head_bg="#831843", soft="#fdf2f8", row_alt="FDF2F8"),
}

# Aliases so natural-language color words map to the canonical palette name
_COLOR_ALIASES = {
    "blue": "azure",
    "dark blue": "navy",
    "light blue": "sky",
    "deep blue": "cobalt",
    "green": "emerald",
    "dark green": "forest",
    "light green": "mint",
    "red": "ruby",
    "dark red": "crimson",
    "light red": "rose",
    "orange": "sunset",
    "dark orange": "ember",
    "light orange": "peach",
    "yellow": "marigold",
    "gold": "brass",
    "purple": "violet",
    "dark purple": "plum",
    "light purple": "lavender",
    "pink": "rose",
    "hot pink": "magenta",
    "black": "mono",
    "dark": "carbon",
    "night": "midnight",
    "grey": "slate",
    "gray": "slate",
    "dark grey": "graphite",
    "dark gray": "graphite",
    "brown": "walnut",
    "tan": "brass",
    "turquoise": "teal",
    "cyan": "ice",
    "aqua": "ice",
}

COLOR_PALETTES = tuple(sorted(_COLOR_PALETTES.keys()))


def _resolve_color_palette(name: str | None) -> dict | None:
    """Return {accent, head_bg, soft, row_alt} for a named color, or None if
    the caller didn't specify one. Handles aliases like 'blue' → 'azure'."""
    if not name:
        return None
    key = name.strip().lower()
    key = _COLOR_ALIASES.get(key, key)
    return _COLOR_PALETTES.get(key)


# Presets grouped by intended document type. When a skill is called with
# `preset="auto"` (the default), we pick from the matching bucket. Lets us
# keep the look appropriate to the document without the model having to
# pick from 60 names.
_PDF_POOLS: dict = {
    "report": [
        "azure",
        "slate",
        "steel",
        "graphite",
        "emerald",
        "forest",
        "teal",
        "indigo",
        "cobalt",
        "navy",
        "oxford",
        "journal",
        "minimal",
        "midnight",
        "obsidian",
        "carbon",
        "lab",
        "blueprint",
        "duotone",
        "academic",
        "cambridge",
        "claret",
        "sage",
    ],
    "letter": [
        "minimal",
        "bone",
        "oxford",
        "navy",
        "journal",
        "whisper",
        "walnut",
        "embassy",
        "manuscript",
        "parchment",
        "academic",
        "slate",
        "graphite",
        "cambridge",
        "broadsheet",
    ],
    "certificate": [
        "embassy",
        "regal",
        "heritage",
        "navy",
        "oxford",
        "walnut",
        "brass",
        "claret",
        "cambridge",
        "parchment",
        "vogue",
        "broadsheet",
        "corporate",
    ],
    "flyer": [
        "sunset",
        "ember",
        "marigold",
        "neon",
        "tropical",
        "sunrise",
        "bubblegum",
        "pastel",
        "citrus",
        "peach",
        "rose",
        "magenta",
        "amber",
        "teal",
        "emerald",
        "indigo",
        "violet",
    ],
    "menu": [
        "editorial",
        "terracotta",
        "walnut",
        "brass",
        "embassy",
        "regal",
        "parchment",
        "rose",
        "amber",
        "sage",
        "manuscript",
        "vogue",
    ],
    "brochure": [
        "azure",
        "ocean",
        "emerald",
        "teal",
        "indigo",
        "violet",
        "sunset",
        "amber",
        "rose",
        "navy",
        "graphite",
        "minimal",
        "marigold",
        "tropical",
        "sage",
    ],
    "newsletter": [
        "editorial",
        "journal",
        "broadsheet",
        "manuscript",
        "vogue",
        "minimal",
        "duotone",
        "graphite",
        "claret",
        "navy",
        "oxford",
    ],
    "card": [
        "minimal",
        "bone",
        "graphite",
        "slate",
        "navy",
        "walnut",
        "azure",
        "carbon",
        "midnight",
        "whisper",
        "duotone",
        "emerald",
        "crimson",
        "violet",
    ],
    "general": [
        "azure",
        "slate",
        "minimal",
        "navy",
        "graphite",
        "emerald",
        "indigo",
        "editorial",
    ],
}

# Keyword hints → preferred preset. Gives the skill a best guess when the
# caller provides title / content / subject text; otherwise we fall back
# to random choice from the skill's pool.
_PDF_KEYWORD_HINTS: list[tuple[tuple[str, ...], str]] = [
    (("wedding", "bridal", "bride", "groom"), "parchment"),
    (("birthday", "party", "celebration", "rsvp"), "bubblegum"),
    (("halloween", "horror", "spooky"), "obsidian"),
    (("christmas", "holiday", "festive"), "crimson"),
    (("finance", "financial", "earnings", "revenue", "q1", "q2", "q3", "q4", "quarterly"), "navy"),
    (("tech", "engineering", "developer", "software", "api", "architecture", "infrastructure"), "carbon"),
    (("research", "paper", "thesis", "study", "academic", "university"), "academic"),
    (("legal", "law firm", "contract", "agreement", "attorney"), "oxford"),
    (("medical", "health", "clinical", "patient"), "teal"),
    (("startup", "pitch", "investor", "fundraising"), "indigo"),
    (("menu", "restaurant", "bistro", "cafe", "café"), "walnut"),
    (("resort", "beach", "vacation", "travel", "tropical"), "tropical"),
    (("summer", "festival", "sunshine", "garden"), "sunset"),
    (("minimalist", "simple", "clean"), "minimal"),
    (("retro", "vintage", "classic"), "terracotta"),
    (("nature", "eco", "sustainability", "environment"), "forest"),
    (("spring", "floral"), "rose"),
    (("government", "official", "ministry"), "embassy"),
    (("music", "concert", "band", "album"), "neon"),
    (("architecture", "construction", "blueprint"), "blueprint"),
]


def _pick_pdf_preset(skill_kind: str, *text_hints: str) -> str:
    """Pick a preset name for `skill_kind` informed by `text_hints` (title,
    content, subject, etc). Walks the keyword hint table first; if nothing
    matches, deterministically hashes the hints to spin the preset pool so
    we don't always return the same one but the same input stays stable.
    """
    import re as _re

    haystack = " ".join(t for t in text_hints if t).lower()
    if haystack:
        # Whole-word tokenization — avoids false positives like "eco" in
        # "recognition" or "q1" in "iraq1984".
        tokens = set(_re.findall(r"[a-z0-9]+", haystack))
        for keywords, preset in _PDF_KEYWORD_HINTS:
            # Each keyword may itself be multi-word (e.g. "law firm"). For
            # multi-word keywords fall back to substring-in-haystack; for
            # single-word keywords require a token match.
            def matches(k: str) -> bool:
                if " " in k:
                    return k in haystack
                return k in tokens

            if any(matches(k) for k in keywords):
                pool = _PDF_POOLS.get(skill_kind) or _PDF_POOLS["general"]
                if preset in pool:
                    return preset
                if preset in _PDF_PRESETS:
                    return preset
    pool = _PDF_POOLS.get(skill_kind) or _PDF_POOLS["general"]
    # Deterministic "random" from hashed input so repeating the same
    # request gives the same preset but different requests vary.
    import hashlib as _hl

    h = int(_hl.sha1(haystack.encode("utf-8")).hexdigest()[:8], 16) if haystack else 0
    return pool[h % len(pool)]


_NAMED_COLORS: dict = {
    "black": "#000000",
    "white": "#ffffff",
    "red": "#dc2626",
    "green": "#16a34a",
    "blue": "#2563eb",
    "yellow": "#eab308",
    "orange": "#ea580c",
    "purple": "#7c3aed",
    "pink": "#ec4899",
    "brown": "#92400e",
    "grey": "#6b7280",
    "gray": "#6b7280",
    "light grey": "#e5e7eb",
    "light gray": "#e5e7eb",
    "dark grey": "#1f2937",
    "dark gray": "#1f2937",
    "cream": "#fef3c7",
    "beige": "#fde68a",
    "light green": "#bbf7d0",
    "mint": "#d1fae5",
    "pale green": "#d1fae5",
    "light blue": "#bfdbfe",
    "sky": "#e0f2fe",
    "light pink": "#fbcfe8",
    "pale pink": "#fce7f3",
    "light yellow": "#fef9c3",
    "pale yellow": "#fef3c7",
    "light purple": "#ddd6fe",
    "lavender": "#e9d5ff",
    "light orange": "#fed7aa",
    "peach": "#ffedd5",
    "light red": "#fecaca",
    "coral": "#fecdd3",
}


def _coerce_color(value: str | None) -> str | None:
    """Accept either a '#rrggbb' / 'rrggbb' string or a natural-language name,
    return a '#rrggbb' string, or None if the input can't be parsed."""
    if not value:
        return None
    s = str(value).strip().lower()
    if s.startswith("#") and len(s) in (4, 7):
        return s
    if len(s) in (3, 6) and all(c in "0123456789abcdef" for c in s):
        return "#" + s
    return _NAMED_COLORS.get(s)


def _pdf_theme(
    name: str,
    skill_kind: str = "general",
    *hints: str,
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Resolve a preset name into a full palette ready for reportlab.

    - `name="auto"` (default): skill picks a preset based on `skill_kind`
      + text `hints`.
    - `color="emerald"` (optional): overrides ONLY the color triple
      (accent / head_bg / soft) while keeping the preset's fonts and
      decoration. Use this for "same look but different color" refinement.
    - `background_color`, `text_color`, `accent_color` (optional): surgical
      overrides that win over both preset and color palette. Accept either
      a named color ("black", "light green", "mint") or a hex ("#09090b").
      Use these when the user explicitly asks for a specific color.
    """
    from reportlab.lib import colors

    key = (name or "auto").strip().lower()
    if key in ("auto", "", "default"):
        key = _pick_pdf_preset(skill_kind, *hints)
    p = _PDF_PRESETS.get(key) or _PDF_PRESETS["azure"]
    # Optional color override — replace accent/head_bg/soft, keep everything
    # else (font family, decoration, page_bg, text color).
    override = _resolve_color_palette(color)
    if override:
        p = {**p, "accent": override["accent"], "head_bg": override["head_bg"], "soft": override["soft"]}
    family = _PDF_FONT_STACKS.get(p.get("family", "modern_sans"))
    text = p.get("text", "#171717")
    muted = p.get("muted", "#6f6f6f")
    page_bg = p.get("page_bg")

    # Surgical overrides — these win over the preset and the color palette.
    bg_hex = _coerce_color(background_color)
    text_hex = _coerce_color(text_color)
    accent_hex = _coerce_color(accent_color)
    if bg_hex:
        page_bg = bg_hex
    if text_hex:
        text = text_hex
    accent = accent_hex if accent_hex else p["accent"]
    head_bg = p["head_bg"]
    # When user sets a background but not an accent, re-derive head_bg
    # from the text color so headers stay legible against the new bg.
    if bg_hex and not accent_hex:
        # Keep the preset's accent if it contrasts with the new bg;
        # otherwise fall back to the text color.
        head_bg = p["head_bg"]

    t = {
        "name": key,
        "color_override": _resolve_color_palette(color) and (color or "").lower(),
        "accent": accent,
        "text": text,
        "muted": muted,
        "soft": p["soft"],
        "head_bg": head_bg,
        "page_bg": page_bg,
        "heading_font": family["heading"],
        "body_font": family["body"],
        "decoration": p["decoration"],
    }
    t["accent_c"] = colors.HexColor(t["accent"])
    t["text_c"] = colors.HexColor(t["text"])
    t["muted_c"] = colors.HexColor(t["muted"])
    t["soft_c"] = colors.HexColor(t["soft"])
    t["head_bg_c"] = colors.HexColor(t["head_bg"])
    t["page_bg_c"] = colors.HexColor(t["page_bg"]) if t["page_bg"] else None
    return t


# Public preset list — exposed so skills.py / UI can enumerate them.
PDF_PRESETS = tuple(sorted(_PDF_PRESETS.keys()))
# Backward-compat alias: the short list of legacy theme names.
PDF_THEMES = (
    "auto",
    "azure",
    "midnight",
    "minimal",
    "navy",
    "editorial",
    "sunset",
    "academic",
)


def generate_pdf_report(
    title: str,
    sections: list[dict] | None = None,
    content: str | None = None,
    filename: str | None = None,
    author: str = "Mio",
    preset: str = "auto",
    color: str | None = None,
    theme: str | None = None,  # legacy alias for preset
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Rich PDF with headings, paragraphs, bullet lists, tables, and charts.

    Two input modes:

    1. `content` (RECOMMENDED): a markdown string. Headings (# / ## / ###),
       paragraphs, bulleted lists, and pipe tables are parsed automatically.
       Easiest for the model — just write markdown.

    2. `sections`: explicit list of block dicts for advanced layouts
       (embedded charts, images, page breaks):
         { "kind": "heading", "text": "...", "level": 1|2|3 }
         { "kind": "paragraph", "text": "..." }
         { "kind": "bullets", "items": ["...", "..."] }
         { "kind": "table", "headers": [...], "rows": [[...],...] }
         { "kind": "chart", "chart_type": "bar|hbar|line|pie",
           "title": "...", "labels": [...], "values": [...] }
         { "kind": "image", "path": "..." }
         { "kind": "pagebreak" }

    You can pass both — `content` is rendered first, then `sections` appended.
    """
    # Merge content (markdown) + sections (structured) into one list.
    merged: list = []
    if content:
        merged.extend(_markdown_to_sections(content))
    if sections:
        if isinstance(sections, (list, tuple)):
            merged.extend(sections)
        else:
            merged.append(sections)
    sections = merged or sections  # keep original list if both were empty
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
            Table,
            TableStyle,
            Image,
            ListFlowable,
            ListItem,
            PageBreak,
        )
    except ImportError:
        return {"skill": "generate_pdf_report", "error": "reportlab not installed."}

    out = _output_path(filename, ".pdf")

    # --- Resolve preset palette + page decoration (auto-picks by default) ---
    T = _pdf_theme(
        preset or theme,
        "report",
        title,
        (content or "")[:400],
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    ACCENT = T["accent_c"]
    MUTED = T["muted_c"]
    TEXT = T["text_c"]
    SOFT_BG = T["soft_c"]
    HEAD_BG = T["head_bg_c"]
    PAGE_BG = T["page_bg_c"]
    HEADING_F = T["heading_font"]
    BODY_F = T["body_font"]
    DECORATION = T["decoration"]

    def _on_page(canvas, doc_):
        canvas.saveState()
        # Page background (dark themes)
        if PAGE_BG is not None:
            canvas.setFillColor(PAGE_BG)
            canvas.rect(0, 0, A4[0], A4[1], fill=1, stroke=0)
        # Top decoration (bar | rule | solid block | none)
        if DECORATION == "bar":
            canvas.setFillColor(ACCENT)
            canvas.rect(0, A4[1] - 6 * mm, A4[0], 6 * mm, fill=1, stroke=0)
        elif DECORATION == "rule":
            canvas.setStrokeColor(ACCENT)
            canvas.setLineWidth(1.2)
            canvas.line(18 * mm, A4[1] - 16 * mm, A4[0] - 18 * mm, A4[1] - 16 * mm)
        elif DECORATION == "block":
            canvas.setFillColor(HEAD_BG)
            canvas.rect(0, A4[1] - 18 * mm, A4[0], 18 * mm, fill=1, stroke=0)
        # Footer
        canvas.setStrokeColor(MUTED)
        canvas.setLineWidth(0.3)
        canvas.line(18 * mm, 14 * mm, A4[0] - 18 * mm, 14 * mm)
        canvas.setFont(BODY_F, 8)
        canvas.setFillColor(MUTED)
        canvas.drawString(18 * mm, 9 * mm, title[:100])
        canvas.drawRightString(A4[0] - 18 * mm, 9 * mm, f"Page {doc_.page}")
        canvas.restoreState()

    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        leftMargin=22 * mm,
        rightMargin=22 * mm,
        topMargin=24 * mm,
        bottomMargin=22 * mm,
        title=title,
        author=author,
    )

    base = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "PMTitle",
        parent=base["Title"],
        fontName=HEADING_F,
        fontSize=26,
        leading=30,
        textColor=TEXT,
        spaceAfter=4,
        alignment=TA_LEFT,
    )
    subtitle_style = ParagraphStyle(
        "PMSubtitle",
        parent=base["Normal"],
        fontName=BODY_F,
        fontSize=10,
        leading=14,
        textColor=MUTED,
        spaceAfter=18,
        alignment=TA_LEFT,
    )
    h1 = ParagraphStyle(
        "PMH1",
        parent=base["Heading1"],
        fontName=HEADING_F,
        fontSize=17,
        leading=22,
        textColor=TEXT,
        spaceBefore=18,
        spaceAfter=6,
        keepWithNext=1,
        borderPadding=(0, 0, 4, 0),
        borderColor=ACCENT,
        borderWidth=0,
    )
    h2 = ParagraphStyle(
        "PMH2",
        parent=base["Heading2"],
        fontName=HEADING_F,
        fontSize=13,
        leading=17,
        textColor=ACCENT,
        spaceBefore=14,
        spaceAfter=4,
        keepWithNext=1,
    )
    h3 = ParagraphStyle(
        "PMH3",
        parent=base["Heading3"],
        fontName=HEADING_F,
        fontSize=11,
        leading=15,
        textColor=TEXT,
        spaceBefore=10,
        spaceAfter=3,
        keepWithNext=1,
    )
    body = ParagraphStyle(
        "PMBody",
        parent=base["BodyText"],
        fontName=BODY_F,
        fontSize=10.5,
        leading=15,
        textColor=TEXT,
        spaceAfter=6,
        alignment=TA_LEFT,
    )
    # Accept `sections` as a list; tolerate accidental single-dict or JSON
    # string the model sometimes produces.
    if isinstance(sections, str):
        import json as _json

        try:
            sections = _json.loads(sections)
        except Exception:
            sections = [{"kind": "paragraph", "text": sections}]
    if isinstance(sections, dict):
        sections = [sections]
    if not isinstance(sections, list) or not sections:
        return {
            "skill": "generate_pdf_report",
            "error": "sections must be a non-empty list of block dicts",
        }

    # Coerce varied section shapes into canonical block dicts so the build
    # doesn't silently drop content when the model emits different formats.
    def _coerce(s):
        if isinstance(s, str):
            if s.strip():
                return [{"kind": "paragraph", "text": s}]
            return []
        if not isinstance(s, dict):
            return []
        if s.get("kind"):
            return [s]
        # No 'kind' — try to infer from keys. Support common variants the
        # model produces when it ignores the schema.
        out = []
        for k in ("heading", "title", "header"):
            if s.get(k):
                out.append({"kind": "heading", "text": str(s[k]), "level": int(s.get("level", 2))})
                break
        for k in ("text", "content", "paragraph", "body", "description"):
            if s.get(k) and isinstance(s[k], str):
                out.append({"kind": "paragraph", "text": s[k]})
                break
        for k in ("bullets", "items", "list", "points"):
            if s.get(k):
                items = s[k]
                if isinstance(items, str):
                    items = [items]
                out.append({"kind": "bullets", "items": list(items)})
                break
        if s.get("rows") or s.get("table"):
            tbl = s.get("table") if isinstance(s.get("table"), dict) else s
            out.append({"kind": "table", "headers": tbl.get("headers") or [], "rows": tbl.get("rows") or []})
        if s.get("chart_type") or s.get("chart"):
            ch = s.get("chart") if isinstance(s.get("chart"), dict) else s
            out.append(
                {
                    "kind": "chart",
                    "chart_type": ch.get("chart_type", "bar"),
                    "title": ch.get("title", ""),
                    "labels": ch.get("labels") or [],
                    "values": ch.get("values") or [],
                }
            )
        if out:
            return out
        # Last resort: render the dict as key/value paragraphs so content
        # never silently vanishes.
        kv = []
        for k, v in s.items():
            if isinstance(v, (str, int, float)):
                kv.append({"kind": "paragraph", "text": f"{k}: {v}"})
        return kv

    normalized: list[dict] = []
    for s in sections:
        normalized.extend(_coerce(s))
    sections = normalized
    if not sections:
        return {
            "skill": "generate_pdf_report",
            "error": "sections contained no renderable blocks",
        }

    import datetime as _dt

    story = [
        Paragraph(title, title_style),
        Paragraph(
            f"{_dt.date.today().strftime('%B %d, %Y')}" + (f" · {author}" if author and author != "Mio" else ""),
            subtitle_style,
        ),
    ]

    import tempfile
    import os

    chart_tmpfiles: list[str] = []

    def _chart_png(sec: dict) -> str | None:
        try:
            import matplotlib

            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return None
        labels = sec.get("labels") or []
        values = sec.get("values") or []
        if not labels or not values or len(labels) != len(values):
            return None
        fig, ax = plt.subplots(figsize=(7.5, 4), dpi=150)
        chart_color = T["accent"]
        ct = (sec.get("chart_type") or "bar").lower()
        if ct == "bar":
            ax.bar(labels, values, color=chart_color)
            plt.setp(ax.get_xticklabels(), rotation=25, ha="right")
        elif ct == "hbar":
            ax.barh(labels, values, color=chart_color)
            ax.invert_yaxis()
        elif ct == "line":
            ax.plot(labels, values, marker="o", color=chart_color, linewidth=2)
            ax.grid(True, alpha=0.3)
        elif ct == "pie":
            ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90)
            ax.axis("equal")
        ax.set_title(sec.get("title", ""), fontsize=12)
        fig.tight_layout()
        f = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(f.name, bbox_inches="tight")
        plt.close(fig)
        chart_tmpfiles.append(f.name)
        return f.name

    for sec in sections:
        kind = (sec.get("kind") or "").lower()
        if kind == "heading":
            lvl = int(sec.get("level", 1))
            style = {1: h1, 2: h2, 3: h3}.get(lvl, h1)
            story.append(Paragraph(sec.get("text", ""), style))
            story.append(Spacer(1, 4))
        elif kind == "paragraph":
            story.append(Paragraph(sec.get("text", ""), body))
            story.append(Spacer(1, 6))
        elif kind == "bullets":
            raw_items = sec.get("items") or []
            if isinstance(raw_items, str):
                raw_items = [raw_items]
            items = [ListItem(Paragraph(str(i), body)) for i in raw_items if i]
            if items:
                story.append(ListFlowable(items, bulletType="bullet", leftIndent=14))
                story.append(Spacer(1, 6))
        elif kind == "table":
            headers = sec.get("headers") or []
            rows = sec.get("rows") or []
            if not headers and not rows:
                continue  # empty table — skip
            if not headers and rows:
                # Infer column count from first row
                headers = [f"Col {i + 1}" for i in range(len(rows[0]))]
            # Normalize all rows to header length
            ncols = len(headers)
            norm_rows = []
            for r in rows:
                r = list(r) if isinstance(r, (list, tuple)) else [r]
                if len(r) < ncols:
                    r = r + [""] * (ncols - len(r))
                else:
                    r = r[:ncols]
                norm_rows.append(r)
            data = [[str(x) for x in headers]] + [[str(x) for x in r] for r in norm_rows]
            tbl = Table(data, repeatRows=1)
            tbl.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), HEAD_BG),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("FONTSIZE", (0, 0), (-1, 0), 9.5),
                        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                        ("FONTSIZE", (0, 1), (-1, -1), 9),
                        ("TEXTCOLOR", (0, 1), (-1, -1), TEXT),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [SOFT_BG, colors.white]),
                        ("LINEBELOW", (0, 0), (-1, 0), 0.6, ACCENT),
                        ("LINEABOVE", (0, 0), (-1, 0), 0, HEAD_BG),
                        ("LINEBELOW", (0, -1), (-1, -1), 0.4, colors.HexColor("#d7dde5")),
                        ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 8),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                        ("TOPPADDING", (0, 0), (-1, -1), 7),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                    ]
                )
            )
            story.append(tbl)
            story.append(Spacer(1, 8))
        elif kind == "chart":
            png = _chart_png(sec)
            if png:
                story.append(Image(png, width=160 * mm, height=88 * mm))
                story.append(Spacer(1, 6))
        elif kind == "image":
            p = sec.get("path")
            if p:
                try:
                    source_path = downloads_input_path(p)
                    with open_binary_no_follow(
                        source_path,
                        max_bytes=16 * 1024 * 1024,
                    ) as source:
                        image_data = io.BytesIO(source.read())
                except (OSError, UnsafePathError) as exc:
                    return {
                        "skill": "generate_pdf_report",
                        "error": f"invalid image input: {exc}",
                    }
                story.append(Image(image_data, width=160 * mm, height=88 * mm))
                story.append(Spacer(1, 6))
        elif kind == "pagebreak":
            story.append(PageBreak())

    try:
        doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    except Exception as e:
        for f in chart_tmpfiles:
            try:
                os.unlink(f)
            except Exception:
                pass
        return {
            "skill": "generate_pdf_report",
            "error": f"reportlab build failed: {type(e).__name__}: {e}",
        }

    for f in chart_tmpfiles:
        try:
            os.unlink(f)
        except Exception:
            pass

    return {
        "skill": "generate_pdf_report",
        "path": str(out),
        "filename": out.name,
    }


# ===========================================================================
# Specialized PDF templates. Each of these takes a small, opinionated schema
# and produces a layout that `generate_pdf_report` can't naturally express
# (multi-column, borders, centered hero, hand-positioned blocks, etc.).
# ===========================================================================


def _rl_imports():
    """Lazy-import reportlab pieces shared by the specialized templates."""
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas as _canvas

    return A4, landscape, mm, colors, _canvas


def generate_letter(
    recipient_name: str,
    recipient_address: str,
    body: str,
    sender_name: str = "",
    sender_address: str = "",
    subject: str = "",
    salutation: str = "",
    closing: str = "Sincerely,",
    date: str = "",
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Formal business letter with letterhead, date, recipient block,
    salutation, body paragraphs, closing, and signature line.

    `body` is plain text; blank lines separate paragraphs. `recipient_address`
    and `sender_address` accept multi-line strings (newline-separated).
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.enums import TA_LEFT, TA_RIGHT
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
        )
    except ImportError:
        return {"skill": "generate_letter", "error": "reportlab not installed."}
    import datetime as _dt

    T = _pdf_theme(
        preset,
        "letter",
        subject,
        (body or "")[:300],
        recipient_name,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    out = _output_path(filename, ".pdf")

    def _on_page(canvas, doc_):
        canvas.saveState()
        if T["page_bg_c"] is not None:
            canvas.setFillColor(T["page_bg_c"])
            canvas.rect(0, 0, A4[0], A4[1], fill=1, stroke=0)
        # Slim rule along the top-right for a light letterhead feel
        canvas.setStrokeColor(T["accent_c"])
        canvas.setLineWidth(1)
        canvas.line(A4[0] - 60 * mm, A4[1] - 18 * mm, A4[0] - 20 * mm, A4[1] - 18 * mm)
        canvas.restoreState()

    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        leftMargin=25 * mm,
        rightMargin=25 * mm,
        topMargin=28 * mm,
        bottomMargin=22 * mm,
        title=subject or f"Letter to {recipient_name}",
    )
    base = getSampleStyleSheet()
    addr = ParagraphStyle(
        "LTAddr", parent=base["Normal"], fontName=T["body_font"], fontSize=10, leading=13.5, textColor=T["text_c"]
    )
    addr_right = ParagraphStyle("LTAddrR", parent=addr, alignment=TA_RIGHT)
    subj_style = ParagraphStyle(
        "LTSubj", parent=addr, fontName=T["heading_font"], fontSize=11.5, spaceBefore=14, spaceAfter=10
    )
    body_style = ParagraphStyle("LTBody", parent=addr, fontSize=10.5, leading=15, spaceAfter=10, alignment=TA_LEFT)
    signature_style = ParagraphStyle("LTSig", parent=addr, fontSize=10.5, leading=14, spaceBefore=40)

    story: list = []
    if sender_address:
        story.append(Paragraph(sender_address.replace("\n", "<br/>"), addr_right))
        story.append(Spacer(1, 18))
    story.append(Paragraph(date or _dt.date.today().strftime("%B %d, %Y"), addr_right))
    story.append(Spacer(1, 16))
    if recipient_name or recipient_address:
        blk = ""
        if recipient_name:
            blk += f"{recipient_name}<br/>"
        if recipient_address:
            blk += recipient_address.replace("\n", "<br/>")
        story.append(Paragraph(blk, addr))
        story.append(Spacer(1, 12))
    if subject:
        story.append(Paragraph(f"<b>Re: {subject}</b>", subj_style))
    story.append(Paragraph(salutation or f"Dear {recipient_name or 'Sir/Madam'},", body_style))
    for para in (body or "").split("\n\n"):
        if para.strip():
            story.append(Paragraph(para.replace("\n", "<br/>"), body_style))
    story.append(Paragraph(closing, body_style))
    if sender_name:
        story.append(Paragraph(f"<br/><br/>_________________________<br/>{sender_name}", signature_style))

    try:
        doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    except Exception as e:
        return {"skill": "generate_letter", "error": f"build failed: {type(e).__name__}: {e}"}
    return {"skill": "generate_letter", "path": str(out), "filename": out.name}


def generate_certificate(
    recipient: str,
    achievement: str,
    issuer: str = "",
    date: str = "",
    signatures: list[dict] | None = None,
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
    orientation: str = "landscape",
) -> dict:
    """Ornate bordered certificate with centered title, recipient, achievement
    text, issue date, and optional signature lines. Default layout is
    landscape A4. `signatures` is a list of {name, role} dicts, placed
    side-by-side at the bottom.
    """
    try:
        A4, landscape, mm, colors, canvas_mod = _rl_imports()
    except ImportError:
        return {"skill": "generate_certificate", "error": "reportlab not installed."}

    T = _pdf_theme(
        preset,
        "certificate",
        achievement,
        recipient,
        issuer,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    out = _output_path(filename, ".pdf")
    page = landscape(A4) if (orientation or "landscape").lower() == "landscape" else A4
    W, H = page

    c = canvas_mod.Canvas(str(out), pagesize=page)
    # Background
    if T["page_bg_c"] is not None:
        c.setFillColor(T["page_bg_c"])
        c.rect(0, 0, W, H, fill=1, stroke=0)
    # Double border
    c.setStrokeColor(T["accent_c"])
    c.setLineWidth(4)
    c.rect(14 * mm, 14 * mm, W - 28 * mm, H - 28 * mm, fill=0, stroke=1)
    c.setLineWidth(1)
    c.rect(19 * mm, 19 * mm, W - 38 * mm, H - 38 * mm, fill=0, stroke=1)
    # Corner ornaments
    for cx, cy in [(26, 26), (W / mm - 26, 26), (26, H / mm - 26), (W / mm - 26, H / mm - 26)]:
        c.setFillColor(T["accent_c"])
        c.circle(cx * mm, cy * mm, 3, fill=1, stroke=0)

    # Header
    c.setFillColor(T["accent_c"])
    c.setFont(T["heading_font"], 18)
    c.drawCentredString(W / 2, H - 38 * mm, "CERTIFICATE OF ACHIEVEMENT")
    c.setFont(T["body_font"], 11)
    c.setFillColor(T["muted_c"])
    c.drawCentredString(W / 2, H - 46 * mm, "This certificate is proudly presented to")

    # Recipient
    c.setFillColor(T["text_c"])
    # Pick a larger display size for short names
    name_size = 44 if len(recipient) < 26 else (36 if len(recipient) < 40 else 28)
    c.setFont(T["heading_font"], name_size)
    c.drawCentredString(W / 2, H / 2 + 12 * mm, recipient)
    # Underline swoosh
    c.setStrokeColor(T["accent_c"])
    c.setLineWidth(1.2)
    c.line(W / 2 - 70 * mm, H / 2 + 6 * mm, W / 2 + 70 * mm, H / 2 + 6 * mm)

    # Achievement text (centered, wrapped)
    c.setFillColor(T["text_c"])
    c.setFont(T["body_font"], 13)
    # Naive wrap — split at word boundaries to ~72 chars per line
    words = (achievement or "").split()
    lines, buf = [], ""
    for w in words:
        if len(buf) + len(w) + 1 > 72:
            lines.append(buf)
            buf = w
        else:
            buf = (buf + " " + w).strip()
    if buf:
        lines.append(buf)
    y = H / 2 - 8 * mm
    for ln in lines[:4]:
        c.drawCentredString(W / 2, y, ln)
        y -= 6 * mm

    # Signatures row
    sigs = signatures or []
    if not sigs and issuer:
        sigs = [{"name": issuer, "role": "Issued by"}]
    if sigs:
        n = len(sigs)
        slot = (W - 60 * mm) / n
        for i, sig in enumerate(sigs):
            x = 30 * mm + slot * i + slot / 2
            c.setStrokeColor(T["muted_c"])
            c.setLineWidth(0.5)
            c.line(x - 28 * mm, 42 * mm, x + 28 * mm, 42 * mm)
            c.setFont(T["heading_font"], 11)
            c.setFillColor(T["text_c"])
            c.drawCentredString(x, 36 * mm, sig.get("name", ""))
            c.setFont(T["body_font"], 9)
            c.setFillColor(T["muted_c"])
            c.drawCentredString(x, 31 * mm, sig.get("role", ""))

    # Date
    import datetime as _dt

    c.setFont(T["body_font"], 10)
    c.setFillColor(T["muted_c"])
    c.drawCentredString(W / 2, 22 * mm, f"Issued on {date or _dt.date.today().strftime('%B %d, %Y')}")
    c.showPage()
    c.save()
    return {"skill": "generate_certificate", "path": str(out), "filename": out.name}


def generate_flyer(
    title: str,
    subtitle: str = "",
    body: str = "",
    call_to_action: str = "",
    footer: str = "",
    image_path: str = "",
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Single-page poster-style flyer: hero title at top, optional hero image
    under it, body copy, a prominent call-to-action pill, and a footer.
    """
    try:
        A4, _, mm, colors, canvas_mod = _rl_imports()
    except ImportError:
        return {"skill": "generate_flyer", "error": "reportlab not installed."}
    T = _pdf_theme(
        preset,
        "flyer",
        title,
        subtitle,
        (body or "")[:300],
        call_to_action,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    hero_image = None
    if image_path:
        try:
            source_path = downloads_input_path(image_path)
            with open_binary_no_follow(
                source_path,
                max_bytes=16 * 1024 * 1024,
            ) as source:
                hero_image = io.BytesIO(source.read())
        except (OSError, UnsafePathError) as exc:
            return {"skill": "generate_flyer", "error": f"invalid image input: {exc}"}
    out = _output_path(filename, ".pdf")
    W, H = A4
    c = canvas_mod.Canvas(str(out), pagesize=A4)

    # Background
    if T["page_bg_c"] is not None:
        c.setFillColor(T["page_bg_c"])
        c.rect(0, 0, W, H, fill=1, stroke=0)
    # Top color block
    c.setFillColor(T["accent_c"])
    c.rect(0, H - 60 * mm, W, 60 * mm, fill=1, stroke=0)
    # Diagonal accent stripe
    c.setFillColor(T["head_bg_c"])
    p = c.beginPath()
    p.moveTo(0, H - 60 * mm)
    p.lineTo(W, H - 30 * mm)
    p.lineTo(W, H - 60 * mm)
    p.close()
    c.drawPath(p, fill=1, stroke=0)

    # Title
    c.setFillColor(colors.white)
    c.setFont(T["heading_font"], 44)
    c.drawString(20 * mm, H - 42 * mm, title[:40])
    if subtitle:
        c.setFont(T["body_font"], 14)
        c.drawString(20 * mm, H - 52 * mm, subtitle[:80])

    # Hero image
    body_top = H - 80 * mm
    if hero_image is not None:
        try:
            from reportlab.lib.utils import ImageReader

            img = ImageReader(hero_image)
            iw, ih = img.getSize()
            target_w = W - 40 * mm
            target_h = target_w * ih / iw
            if target_h > 100 * mm:
                target_h = 100 * mm
                target_w = target_h * iw / ih
            c.drawImage(img, (W - target_w) / 2, body_top - target_h, width=target_w, height=target_h, mask="auto")
            body_top = body_top - target_h - 10 * mm
        except Exception:
            pass

    # Body
    c.setFillColor(T["text_c"])
    c.setFont(T["body_font"], 12)
    lines, buf = [], ""
    for w in (body or "").split():
        if len(buf) + len(w) + 1 > 80:
            lines.append(buf)
            buf = w
        else:
            buf = (buf + " " + w).strip()
    if buf:
        lines.append(buf)
    y = body_top - 6 * mm
    for ln in lines[:12]:
        c.drawString(20 * mm, y, ln)
        y -= 6 * mm

    # Call to action pill
    if call_to_action:
        pill_w = min(140 * mm, 8 * len(call_to_action) + 40)
        pill_h = 16 * mm
        cx = W / 2
        cy = 40 * mm
        c.setFillColor(T["accent_c"])
        c.roundRect(cx - pill_w / 2, cy - pill_h / 2, pill_w, pill_h, radius=pill_h / 2, fill=1, stroke=0)
        c.setFillColor(colors.white)
        c.setFont(T["heading_font"], 14)
        c.drawCentredString(cx, cy - 2, call_to_action[:60])

    # Footer
    if footer:
        c.setFillColor(T["muted_c"])
        c.setFont(T["body_font"], 9)
        c.drawCentredString(W / 2, 16 * mm, footer[:140])

    c.showPage()
    c.save()
    return {"skill": "generate_flyer", "path": str(out), "filename": out.name}


def generate_menu(
    restaurant_name: str,
    tagline: str = "",
    sections: list[dict] | None = None,
    footer: str = "",
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Restaurant menu: big restaurant name, tagline, then two-column grid of
    category sections. `sections` is a list of:
        { "name": "Appetizers",
          "items": [
            {"name": "Caprese", "description": "tomato, mozzarella…", "price": "$12"},
            ...
          ] }
    """
    try:
        A4, _, mm, colors, canvas_mod = _rl_imports()
        from reportlab.pdfbase.pdfmetrics import stringWidth
    except ImportError:
        return {"skill": "generate_menu", "error": "reportlab not installed."}

    T = _pdf_theme(
        preset,
        "menu",
        restaurant_name,
        tagline,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    out = _output_path(filename, ".pdf")
    W, H = A4
    c = canvas_mod.Canvas(str(out), pagesize=A4)
    if T["page_bg_c"] is not None:
        c.setFillColor(T["page_bg_c"])
        c.rect(0, 0, W, H, fill=1, stroke=0)

    # Header
    c.setFillColor(T["text_c"])
    c.setFont(T["heading_font"], 34)
    c.drawCentredString(W / 2, H - 24 * mm, restaurant_name)
    if tagline:
        c.setFont(T["body_font"], 11)
        c.setFillColor(T["muted_c"])
        c.drawCentredString(W / 2, H - 31 * mm, tagline)
    c.setStrokeColor(T["accent_c"])
    c.setLineWidth(1.2)
    c.line(W / 2 - 40 * mm, H - 37 * mm, W / 2 + 40 * mm, H - 37 * mm)

    # Two-column layout
    col_w = (W - 60 * mm) / 2
    col_x = [20 * mm, 20 * mm + col_w + 20 * mm]
    col_y = [H - 52 * mm, H - 52 * mm]
    col = 0
    for sec in sections or []:
        name = sec.get("name", "")
        items = sec.get("items", []) or []
        needed = 14 + 16 * len(items)  # approx height in mm
        if col_y[col] - needed * mm < 30 * mm:
            col = 1 - col
            if col_y[col] - needed * mm < 30 * mm:
                break  # page full
        # Section heading
        c.setFillColor(T["accent_c"])
        c.setFont(T["heading_font"], 14)
        c.drawString(col_x[col], col_y[col], name.upper())
        col_y[col] -= 4 * mm
        c.setStrokeColor(T["muted_c"])
        c.setLineWidth(0.3)
        c.line(col_x[col], col_y[col], col_x[col] + col_w, col_y[col])
        col_y[col] -= 6 * mm
        # Items
        for it in items:
            nm = it.get("name", "")
            desc = it.get("description", "")
            price = it.get("price", "")
            c.setFillColor(T["text_c"])
            c.setFont(T["heading_font"], 11)
            c.drawString(col_x[col], col_y[col], nm)
            if price:
                c.drawRightString(col_x[col] + col_w, col_y[col], str(price))
            col_y[col] -= 4 * mm
            if desc:
                c.setFillColor(T["muted_c"])
                c.setFont(T["body_font"], 9)
                # truncate long descriptions to fit column
                maxw = col_w
                while stringWidth(desc, T["body_font"], 9) > maxw and len(desc) > 10:
                    desc = desc[:-2]
                c.drawString(col_x[col], col_y[col], desc)
                col_y[col] -= 4 * mm
            col_y[col] -= 3 * mm

    if footer:
        c.setFillColor(T["muted_c"])
        c.setFont(T["body_font"], 9)
        c.drawCentredString(W / 2, 14 * mm, footer)

    c.showPage()
    c.save()
    return {"skill": "generate_menu", "path": str(out), "filename": out.name}


def generate_brochure(
    title: str,
    panels: list[dict],
    footer: str = "",
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Landscape tri-fold brochure (3 equal panels). `panels` is a list of
    exactly three dicts: { "heading": "...", "body": "...", "bullets": [...] }.
    Front panel usually goes on the right-most fold.
    """
    try:
        A4, landscape, mm, colors, canvas_mod = _rl_imports()
        from reportlab.pdfbase.pdfmetrics import stringWidth
    except ImportError:
        return {"skill": "generate_brochure", "error": "reportlab not installed."}

    panel_hints = " ".join((p.get("heading", "") + " " + p.get("body", ""))[:200] for p in (panels or [])[:3])
    T = _pdf_theme(
        preset,
        "brochure",
        title,
        panel_hints,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    out = _output_path(filename, ".pdf")
    page = landscape(A4)
    W, H = page
    c = canvas_mod.Canvas(str(out), pagesize=page)
    if T["page_bg_c"] is not None:
        c.setFillColor(T["page_bg_c"])
        c.rect(0, 0, W, H, fill=1, stroke=0)

    panels = (panels or [])[:3]
    while len(panels) < 3:
        panels.append({})

    panel_w = (W - 40 * mm) / 3
    fold_xs = [15 * mm + panel_w, 15 * mm + panel_w * 2 + 5 * mm]
    # Fold guides (very light)
    c.setStrokeColor(T["muted_c"])
    c.setDash(2, 3)
    for fx in fold_xs:
        c.line(fx, 10 * mm, fx, H - 10 * mm)
    c.setDash()

    for i, panel in enumerate(panels):
        x = 20 * mm + i * (panel_w + 5 * mm)
        heading = panel.get("heading", "")
        body = panel.get("body", "")
        bullets = panel.get("bullets", []) or []
        # Panel header bar
        c.setFillColor(T["accent_c"])
        c.rect(x, H - 28 * mm, panel_w - 5 * mm, 4, fill=1, stroke=0)
        c.setFillColor(T["text_c"])
        c.setFont(T["heading_font"], 16)
        c.drawString(x, H - 36 * mm, heading[:40])
        # Body
        c.setFillColor(T["text_c"])
        c.setFont(T["body_font"], 10)
        y = H - 44 * mm
        for para in (body or "").split("\n\n"):
            if not para.strip():
                continue
            words = para.split()
            lines, buf = [], ""
            maxw = panel_w - 10 * mm
            for w in words:
                test = (buf + " " + w).strip()
                if stringWidth(test, T["body_font"], 10) > maxw:
                    lines.append(buf)
                    buf = w
                else:
                    buf = test
            if buf:
                lines.append(buf)
            for ln in lines:
                if y < 30 * mm:
                    break
                c.drawString(x, y, ln)
                y -= 4.6 * mm
            y -= 3 * mm
        # Bullets
        for b in bullets:
            if y < 30 * mm:
                break
            c.setFillColor(T["accent_c"])
            c.circle(x + 1.6, y + 1.5, 0.9, fill=1, stroke=0)
            c.setFillColor(T["text_c"])
            c.setFont(T["body_font"], 10)
            # Truncate long bullet to fit
            s = str(b)
            maxw = panel_w - 12 * mm
            while stringWidth(s, T["body_font"], 10) > maxw and len(s) > 8:
                s = s[:-2] + "…"
            c.drawString(x + 5 * mm, y, s)
            y -= 5.5 * mm

    # Title footer
    c.setFillColor(T["muted_c"])
    c.setFont(T["body_font"], 8)
    c.drawString(20 * mm, 8 * mm, title[:80])
    if footer:
        c.drawRightString(W - 20 * mm, 8 * mm, footer[:80])

    c.showPage()
    c.save()
    return {"skill": "generate_brochure", "path": str(out), "filename": out.name}


def generate_newsletter(
    title: str,
    issue: str = "",
    lead_headline: str = "",
    lead_body: str = "",
    articles: list[dict] | None = None,
    footer: str = "",
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Newsletter-style PDF: masthead, lead story (full-width headline +
    paragraph), then two-column grid of article blocks (heading + body).
    """
    try:
        A4, _, mm, colors, canvas_mod = _rl_imports()
        from reportlab.pdfbase.pdfmetrics import stringWidth
    except ImportError:
        return {"skill": "generate_newsletter", "error": "reportlab not installed."}

    art_hints = " ".join((a.get("heading", "") + " " + a.get("body", ""))[:200] for a in (articles or [])[:4])
    T = _pdf_theme(
        preset,
        "newsletter",
        title,
        lead_headline,
        lead_body[:300] if lead_body else "",
        art_hints,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    out = _output_path(filename, ".pdf")
    W, H = A4
    c = canvas_mod.Canvas(str(out), pagesize=A4)
    if T["page_bg_c"] is not None:
        c.setFillColor(T["page_bg_c"])
        c.rect(0, 0, W, H, fill=1, stroke=0)

    # Masthead
    c.setFillColor(T["accent_c"])
    c.rect(0, H - 8 * mm, W, 3, fill=1, stroke=0)
    c.setFillColor(T["text_c"])
    c.setFont(T["heading_font"], 28)
    c.drawString(18 * mm, H - 22 * mm, title[:50])
    c.setFillColor(T["muted_c"])
    c.setFont(T["body_font"], 10)
    import datetime as _dt

    c.drawRightString(W - 18 * mm, H - 22 * mm, issue or _dt.date.today().strftime("%B %Y"))
    c.setStrokeColor(T["text_c"])
    c.setLineWidth(0.8)
    c.line(18 * mm, H - 28 * mm, W - 18 * mm, H - 28 * mm)

    y_cursor = H - 36 * mm

    # Lead story
    if lead_headline:
        c.setFillColor(T["text_c"])
        c.setFont(T["heading_font"], 20)
        c.drawString(18 * mm, y_cursor, lead_headline[:80])
        y_cursor -= 8 * mm
    if lead_body:
        c.setFillColor(T["text_c"])
        c.setFont(T["body_font"], 11)
        words = lead_body.split()
        lines, buf = [], ""
        maxw = W - 36 * mm
        for w in words:
            test = (buf + " " + w).strip()
            if stringWidth(test, T["body_font"], 11) > maxw:
                lines.append(buf)
                buf = w
            else:
                buf = test
        if buf:
            lines.append(buf)
        for ln in lines[:6]:
            c.drawString(18 * mm, y_cursor, ln)
            y_cursor -= 5 * mm
        y_cursor -= 4 * mm

    # Articles — two columns
    col_w = (W - 46 * mm) / 2
    col_x = [18 * mm, 18 * mm + col_w + 10 * mm]
    col_y = [y_cursor, y_cursor]
    col = 0
    for art in articles or []:
        heading = art.get("heading") or art.get("title") or ""
        body = art.get("body") or art.get("text") or ""
        # Rough vertical estimate
        est = 10 + (len(body) / 60) * 5
        if col_y[col] - est * mm < 25 * mm:
            col = 1 - col
            if col_y[col] - est * mm < 25 * mm:
                break
        c.setFillColor(T["accent_c"])
        c.setFont(T["heading_font"], 12)
        c.drawString(col_x[col], col_y[col], heading[:50])
        col_y[col] -= 6 * mm
        c.setFillColor(T["text_c"])
        c.setFont(T["body_font"], 10)
        words = body.split()
        lines, buf = [], ""
        for w in words:
            test = (buf + " " + w).strip()
            if stringWidth(test, T["body_font"], 10) > col_w:
                lines.append(buf)
                buf = w
            else:
                buf = test
        if buf:
            lines.append(buf)
        for ln in lines:
            if col_y[col] < 25 * mm:
                break
            c.drawString(col_x[col], col_y[col], ln)
            col_y[col] -= 4.5 * mm
        col_y[col] -= 6 * mm

    if footer:
        c.setFillColor(T["muted_c"])
        c.setFont(T["body_font"], 8)
        c.drawCentredString(W / 2, 12 * mm, footer)

    c.showPage()
    c.save()
    return {"skill": "generate_newsletter", "path": str(out), "filename": out.name}


def generate_business_card(
    name: str,
    role: str = "",
    company: str = "",
    email: str = "",
    phone: str = "",
    website: str = "",
    address: str = "",
    filename: str | None = None,
    preset: str = "auto",
    color: str | None = None,
    background_color: str | None = None,
    text_color: str | None = None,
    accent_color: str | None = None,
) -> dict:
    """Single business card at standard US 3.5"×2" size. Centered on an A4
    page for easy printing; laid out with name, role, and contact block.
    """
    try:
        A4, _, mm, colors, canvas_mod = _rl_imports()
    except ImportError:
        return {"skill": "generate_business_card", "error": "reportlab not installed."}

    T = _pdf_theme(
        preset,
        "card",
        name,
        role,
        company,
        color=color,
        background_color=background_color,
        text_color=text_color,
        accent_color=accent_color,
    )
    out = _output_path(filename, ".pdf")
    W, H = A4
    card_w, card_h = 88.9 * mm, 50.8 * mm  # 3.5 x 2 inches
    cx, cy = (W - card_w) / 2, (H - card_h) / 2
    c = canvas_mod.Canvas(str(out), pagesize=A4)
    if T["page_bg_c"] is not None:
        c.setFillColor(T["page_bg_c"])
        c.rect(0, 0, W, H, fill=1, stroke=0)
    # Card background + accent strip on left
    c.setFillColor(colors.white if T["page_bg_c"] is None else T["soft_c"])
    c.setStrokeColor(T["muted_c"])
    c.setLineWidth(0.3)
    c.rect(cx, cy, card_w, card_h, fill=1, stroke=1)
    c.setFillColor(T["accent_c"])
    c.rect(cx, cy, 4 * mm, card_h, fill=1, stroke=0)

    # Name + role
    c.setFillColor(T["text_c"])
    c.setFont(T["heading_font"], 16)
    c.drawString(cx + 10 * mm, cy + card_h - 14 * mm, name[:28])
    if role:
        c.setFont(T["body_font"], 10)
        c.setFillColor(T["muted_c"])
        c.drawString(cx + 10 * mm, cy + card_h - 20 * mm, role[:40])
    if company:
        c.setFont(T["heading_font"], 10)
        c.setFillColor(T["accent_c"])
        c.drawString(cx + 10 * mm, cy + card_h - 26 * mm, company[:40])

    # Contact block
    y = cy + 18 * mm
    c.setFillColor(T["text_c"])
    c.setFont(T["body_font"], 9)
    for label, val in [("E", email), ("T", phone), ("W", website), ("A", address)]:
        if not val:
            continue
        c.setFillColor(T["accent_c"])
        c.drawString(cx + 10 * mm, y, label)
        c.setFillColor(T["text_c"])
        c.drawString(cx + 14 * mm, y, str(val)[:44])
        y -= 4 * mm

    c.showPage()
    c.save()
    return {"skill": "generate_business_card", "path": str(out), "filename": out.name}
